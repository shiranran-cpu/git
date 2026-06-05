from pptx import Presentation

class PPTGenerator:
    def generate(self, pdf_path):
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = 'AI Research PPT'
        out = pdf_path + '.pptx'
        prs.save(out)
        return out
